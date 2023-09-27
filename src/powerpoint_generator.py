from pptx import Presentation
import os
from pathlib import Path

# Function to let the user navigate and select a .txt file
def select_input_file():
    dir_path = './ppt'
    txt_files = [f for f in os.listdir(dir_path) if os.path.isfile(os.path.join(dir_path, f)) and f.endswith('.txt')]

    if not txt_files:
        print("\nNo text files found in the './ppt' directory.")
        return None

    txt_files.sort()  # Sort the list of text files alphabetically

    print("\nWhich txt file do you want to convert to a powerpoint?")
    for idx, file in enumerate(txt_files):
        print(f"{idx + 1}. {file}")

    while True:
        try:
            selected_index = int(input("Enter the number of your selected file: ")) - 1
            if 0 <= selected_index < len(txt_files):
                return os.path.join(dir_path, txt_files[selected_index])
            else:
                print("Invalid choice. Please select a number from the list.")
        except ValueError:
            print("Please enter a valid number.")

# Function to extract slide data from text file
def extract_slide_data(file_name):
    with open(file_name, 'r') as f:
        content = f.read()

    # Filter out lines starting with ##
    content = '\n'.join([line for line in content.split('\n') if not line.strip().startswith('##')])

    # Extract slides
    slides = [s.strip() for s in content.split('--- Slide ---') if s.strip() and '--- End Slide ---' in s]

    slide_data = []
    for slide in slides:
        data = {}
        fields_content = slide.split('--- End Slide ---')[0].strip()
        
        for field in ['Layout', 'Title', 'Content', 'Speaker Note']:
            if f"{field}:" in fields_content:
                start_idx = fields_content.find(f"{field}:") + len(f"{field}:")
                next_field_idx = [fields_content.find(f"{f}:") for f in ['Layout', 'Title', 'Content', 'Speaker Note'] if f"{f}:" in fields_content and fields_content.find(f"{f}:") > start_idx]
                end_idx = min(next_field_idx) if next_field_idx else None
                data[field] = fields_content[start_idx:end_idx].strip()

        # Sanitize the Layout field, only keep digits
        if 'Layout' in data:
            data['Layout'] = ''.join(filter(str.isdigit, data['Layout']))
            
        slide_data.append(data)

    return slide_data


def generate_ppt_from_file(file_path):
    # Extract slide data from input file
    slide_data = extract_slide_data(file_path)

    # Create a new PowerPoint presentation
    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "blank_powerpoint_template.pptx")
    try:
        prs = Presentation(template_path)
    except:
        prs = Presentation()

    # Delete existing slides
    for i in range(len(prs.slides)-1, -1, -1):
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])

    # Generate slides based on slide data
    for slide_info in slide_data:
        slide_layout = int(slide_info.get('Layout', 1))
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
        
        # Set title, content, and speaker note
        if 'Title' in slide_info:
            slide.shapes.title.text = slide_info['Title']
        
        if 'Content' in slide_info:
            slide.placeholders[1].text = slide_info['Content']
        
        if 'Speaker Note' in slide_info:
            slide.notes_slide.notes_text_frame.text = slide_info['Speaker Note']

    # Derive the PowerPoint filename from the input text file
    ppt_filename = os.path.splitext(os.path.basename(file_path))[0] + ".pptx"
    prs.save(os.path.join('ppt', ppt_filename))
    print(f"\nPresentation saved as {ppt_filename}")

def main():
    run_program = True
    while run_program:
        # Get the input file
        file_path = select_input_file()

        if file_path:
            generate_ppt_from_file(file_path)
        else:
            print("No file selected.")
            break

        # Prompt the user to run the program again
        run_again = input("\nDo you want to run the program again? (y/n): ")
        if run_again.lower() == 'n':
            print("\n***Exiting program***\n")
            run_program = False

if __name__ == "__main__":
    main()